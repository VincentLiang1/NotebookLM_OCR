"""Emit the DeckEdit-style PPTX: full-bleed background picture per slide plus
solid-fill cover shapes carrying editable, style-matched text."""
from __future__ import annotations

from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

from .models import ALIGN_CENTER, ALIGN_RIGHT, TextBlock

# ALL latin characters render in a dedicated latin face (user request);
# OOXML fonts are per-character-class, so one run carries <a:latin>=Arial
# for its English and <a:ea>=YaHei for its CJK — no run splitting needed.
# Page titles (>= NARROW_MIN_PT) use Arial Narrow instead (user request:
# the source deck's latin is narrower than Arial and long mixed titles
# like p4 Layer 0 ran visibly long; style.py measures widths to match).
LATIN_FONT = "Arial"
LATIN_FONT_NARROW = "Arial Narrow"


def _latin_font(pt: float) -> str:
    from .style import NARROW_MIN_PT
    return LATIN_FONT_NARROW if pt >= NARROW_MIN_PT else LATIN_FONT

SLIDE_W_EMU = 12192000  # 13.333 in, matches the example deck
EMU_PER_PT = 12700
COVER_PAD_PX = 3     # expand cover box to hide anti-aliased fringe of raster text
LEADING_COMP = 0.20  # em: gap between a top-anchored frame's top and glyph ink
WARP_DROP_RATIO = 0.13  # arch warp: rendered drop per unit of frame height

PP_ALIGN_MAP = {ALIGN_CENTER: PP_ALIGN.CENTER, ALIGN_RIGHT: PP_ALIGN.RIGHT}


class DeckBuilder:
    def __init__(self, page_w_pt: float, page_h_pt: float, font_name: str,
                 cover: bool = True):
        self.prs = Presentation()
        self.prs.slide_width = Emu(SLIDE_W_EMU)
        self.slide_h_emu = round(SLIDE_W_EMU * page_h_pt / page_w_pt)
        self.prs.slide_height = Emu(self.slide_h_emu)
        self.blank_layout = self.prs.slide_layouts[6]
        self.font_name = font_name
        self.cover = cover

    def add_slide(self, png_bytes: bytes, blocks: list[TextBlock],
                  img_w: int, img_h: int,
                  wipes: list[tuple[tuple, tuple]] | None = None,
                  img=None) -> None:
        """wipes: [(bbox_px, rgb)] — text-less cover rectangles that blank
        out regions (e.g. the NotebookLM watermark) with the background
        color. img: the page render (numpy RGB), used to clamp arc cover
        strips to their ribbon."""
        slide = self.prs.slides.add_slide(self.blank_layout)
        pic = slide.shapes.add_picture(
            BytesIO(png_bytes), 0, 0,
            self.prs.slide_width, self.prs.slide_height,
        )
        pic.name = "Image 0"

        def ex(px: float) -> int:
            return round(px * SLIDE_W_EMU / img_w)

        def ey(px: float) -> int:
            return round(px * self.slide_h_emu / img_h)

        for i, (bbox, rgb) in enumerate(wipes or []):
            x0, y0, x1, y1 = bbox
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Emu(max(0, ex(x0))), Emu(max(0, ey(y0))),
                Emu(ex(x1) - ex(x0)), Emu(ey(y1) - ey(y0)),
            )
            shape.name = f"Wipe {i}"
            shape.shadow.inherit = False
            shape.line.fill.background()
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*rgb)

        # same-row OCR boxes can overlap horizontally (the detector pads a
        # box into its neighbor's glyphs); the later-drawn cover then
        # paints over the neighbor's already-drawn text. Trim both boxes
        # to the blank ink gap between them, measured on the render.
        if img is not None:
            self._trim_row_overlaps(blocks, img)

        # arc covers go in first: two arc lines on the same ribbon
        # interpenetrate, and a later line's cover strips must not paint
        # over an earlier line's text segments
        arc_blocks = [(i, b) for i, b in enumerate(blocks)
                      if len(b.lines) == 1 and b.lines[0].arc_sagitta]
        for i, block in arc_blocks:
            self._add_arc_cover(slide, block, i, ex, ey,
                                px_per_pt=img_w / 960.0, img=img)

        for i, block in enumerate(blocks):
            if len(block.lines) == 1 and block.lines[0].arc_sagitta:
                self._add_arc_segments(slide, block, i, ex, ey,
                                       px_per_pt=img_w / 960.0, img=img)
                continue
            vertical = len(block.lines) == 1 and block.style.vertical
            nudge = (0 if vertical
                     else round(LEADING_COMP * block.style.font_pt * EMU_PER_PT))
            tilted = (len(block.lines) == 1 and block.lines[0].angle
                      and block.lines[0].center and block.lines[0].size)
            if tilted:
                # rotated shape: deskewed quad size centered on the quad
                # center; PowerPoint rotates around the shape center, and
                # both angles are clockwise-positive
                ln = block.lines[0]
                width = ex(ln.size[0] + 2 * COVER_PAD_PX)
                height = ey(ln.size[1] + 2 * COVER_PAD_PX)
                left = ex(ln.center[0]) - width // 2
                top = ey(ln.center[1]) - height // 2 - nudge
                height += nudge
            else:
                x0, y0, x1, y1 = block.bbox
                # start/end the cover at the text, not the box edge, when
                # the box overhangs an adjacent graphic (p13 red ✗)
                if len(block.lines) == 1 and block.style.cover_x0_px is not None:
                    x0 = max(x0, block.style.cover_x0_px)
                if len(block.lines) == 1 and block.style.cover_x1_px is not None:
                    x1 = min(x1, block.style.cover_x1_px)
                left = ex(x0 - COVER_PAD_PX)
                width = ex(x1 + COVER_PAD_PX) - left
                # cover height follows the glyph ink band, not the OCR box:
                # detector boxes carry large vertical slack that would paint
                # over diagram lines above/below the text. The text frame
                # top is decoupled from the cover top via a margin inset.
                if len(block.lines) == 1 and block.style.highlight_removed:
                    # a dropped inline highlight box spans (and slightly
                    # overhangs) the OCR box; cover generously past it so no
                    # source fill leaks past the glyph band (user: "remove
                    # it cleanly, no leak")
                    over = max(8.0, 0.12 * (y1 - y0))
                    cov_y0, cov_y1 = y0 - over, y1 + over
                elif len(block.lines) == 1 and block.style.ink_bottom_px:
                    cov_h = block.style.ink_bottom_px - block.style.ink_top_px
                    pad_v = max(4.0, 0.08 * cov_h)
                    cov_y0 = block.style.ink_top_px - pad_v
                    cov_y1 = block.style.ink_bottom_px + pad_v
                else:
                    cov_y0, cov_y1 = y0 - COVER_PAD_PX, y1 + COVER_PAD_PX
                ink_top = ey(block.style.ink_top_px)
                text_top = ink_top - nudge
                top = min(text_top, ey(cov_y0))
                height = ey(cov_y1) - top
                margin_top = max(0, text_top - top)

            fill = self.cover and block.style.bg_rgb is not None
            bg_segs = block.style.bg_segments
            if not tilted and self.cover and bg_segs and len(bg_segs) >= 2:
                # two-tone banner: a FULL-WIDTH base cover in the rightmost
                # fill color (so the seam can never expose the raster even
                # if a cover is nudged), then each earlier segment painted
                # on top up to its boundary. The transparent text box on top
                # carries per-segment color runs (white on the dark fill,
                # dark on the light fill).
                cv_top, cv_h = Emu(max(0, ey(cov_y0))), Emu(ey(cov_y1) - ey(cov_y0))
                full_l = ex(block.bbox[0] - COVER_PAD_PX)
                full_r = ex(block.bbox[2] + COVER_PAD_PX)
                covers = [(full_l, full_r, bg_segs[-1][2])]
                for si in range(len(bg_segs) - 1):
                    sx0, sx1, sbg = bg_segs[si]
                    cl = ex(sx0 - COVER_PAD_PX) if si == 0 else ex(sx0)
                    covers.append((cl, ex(sx1), sbg))
                for si, (cl, cr, sbg) in enumerate(covers):
                    cov = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE, Emu(max(0, cl)), cv_top,
                        Emu(cr - cl), cv_h,
                    )
                    cov.name = f"Text {i} bg{si}"
                    cov.shadow.inherit = False
                    cov.line.fill.background()
                    cov.fill.solid()
                    cov.fill.fore_color.rgb = RGBColor(*sbg)
                fill = False
                top = text_top
                height = ey(cov_y1) - top
                margin_top = 0
            elif not tilted and fill and text_top < ey(cov_y0):
                # the leading-compensation zone above the ink would carry
                # the fill onto the previous line's descenders in tight
                # rows; split into a cover rect plus a transparent text box
                cover = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Emu(max(0, left)),
                    Emu(max(0, ey(cov_y0))),
                    Emu(width), Emu(ey(cov_y1) - ey(cov_y0)),
                )
                cover.name = f"Text {i} bg"
                cover.shadow.inherit = False
                cover.line.fill.background()
                cover.fill.solid()
                cover.fill.fore_color.rgb = RGBColor(*block.style.bg_rgb)
                fill = False
                top = text_top
                height = ey(cov_y1) - top
                margin_top = 0

            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Emu(max(0, left)), Emu(max(0, top)),
                Emu(width), Emu(height),
            )
            if tilted:
                shape.rotation = block.lines[0].angle
                margin_top = 0
            shape.name = f"Text {i}"
            shape.shadow.inherit = False
            shape.line.fill.background()
            if fill:
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(*block.style.bg_rgb)
            else:
                shape.fill.background()

            tf = shape.text_frame
            # single-line shapes must never wrap: a size estimate a hair too
            # large would otherwise break the line and wreck the layout
            tf.word_wrap = len(block.lines) > 1
            tf.auto_size = None
            tf.vertical_anchor = MSO_ANCHOR.TOP
            tf.margin_left = tf.margin_right = tf.margin_bottom = 0
            tf.margin_top = Emu(margin_top)
            if vertical:
                # east-asian vertical text: glyphs stacked top-to-bottom,
                # centered in the column, no leading inset
                tf.word_wrap = False
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                tf.margin_top = 0
                tf._txBody.bodyPr.set("vert", "eaVert")

            for j, line in enumerate(block.lines):
                para = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                para.alignment = (PP_ALIGN.CENTER if vertical
                                  else PP_ALIGN_MAP.get(block.align, PP_ALIGN.LEFT))
                pieces = None
                if len(block.lines) == 1 and block.style.runs:
                    pieces = _split_text_runs(line.text, block.style.runs)
                if pieces is None:
                    pieces = [(line.text, block.style.text_rgb)]
                tail = block.style.superscript_tail if len(block.lines) == 1 else 0
                pieces = _mark_superscript(pieces, tail)
                for piece, rgb, sup in pieces:
                    run = para.add_run()
                    run.text = piece
                    font = run.font
                    font.size = Pt(block.style.font_pt * (0.65 if sup else 1.0))
                    font.bold = block.style.bold
                    font.name = _latin_font(block.style.font_pt)  # <a:latin>
                    font.color.rgb = RGBColor(*rgb)
                    rPr = run._r.get_or_add_rPr()
                    if block.style.strikethrough:
                        rPr.set("strike", "sngStrike")
                    if sup:                       # raised footnote marker
                        rPr.set("baseline", "30000")
                    _set_east_asian_font(run, self.font_name)

    @staticmethod
    def _trim_row_overlaps(blocks, img) -> None:
        """Split horizontally overlapping same-row boxes at the blank
        column gap between their glyphs (page 3 '…書）' / '再批次…': the
        right box starts 43px inside the left one, so the left cover
        painted over the right line's 再). Falls back to the overlap
        midpoint when the seam is solid ink."""
        import numpy as np

        plain = [b for b in blocks
                 if not (len(b.lines) == 1
                         and (b.lines[0].arc_sagitta or b.lines[0].angle))]
        plain.sort(key=lambda b: b.bbox[0])
        for ai, a in enumerate(plain):
            for b in plain[ai + 1:]:
                ax0, ay0, ax1, ay1 = a.bbox
                bx0, by0, bx1, by1 = b.bbox
                if bx0 >= ax1 - COVER_PAD_PX:
                    continue
                oy = min(ay1, by1) - max(ay0, by0)
                if oy < 0.6 * min(ay1 - ay0, by1 - by0):
                    continue
                y_lo, y_hi = int(max(ay0, by0)), int(min(ay1, by1))
                x_lo = max(0, int(bx0 - 8))
                x_hi = min(img.shape[1], int(ax1 + 16))
                boundary = (ax1 + bx0) / 2
                bg = a.style.bg_rgb or b.style.bg_rgb
                if bg is not None and y_hi - y_lo >= 4 and x_hi - x_lo >= 8:
                    win = img[y_lo:y_hi, x_lo:x_hi].astype(int)
                    has_ink = (np.abs(win - np.asarray(bg, dtype=int))
                               .max(axis=2) > 60).sum(axis=0) >= 2
                    runs, s = [], None
                    for j, v in enumerate(has_ink):
                        if not v and s is None:
                            s = j
                        elif v and s is not None:
                            runs.append((s, j))
                            s = None
                    if s is not None:
                        runs.append((s, len(has_ink)))
                    runs = [r for r in runs if r[1] - r[0] >= 4]
                    if runs:
                        lo, hi = max(runs, key=lambda r: r[1] - r[0])
                        boundary = x_lo + (lo + hi) / 2
                a._bbox = (ax0, ay0,
                           min(ax1, boundary - COVER_PAD_PX), ay1)
                b._bbox = (max(bx0, boundary + COVER_PAD_PX), by0, bx1, by1)

    @staticmethod
    def _arc_geometry(block, px_per_pt: float, img=None):
        ln = block.lines[0]
        x0, y0, x1, y1 = ln.bbox
        glyph_h = min(block.style.font_pt * 1.2 * px_per_pt, (y1 - y0) * 0.8)
        if ln.arc_sagitta > 0:  # arch up: middle at top, edges at bottom
            y_mid, y_edge = y0 + glyph_h / 2, y1 - glyph_h / 2
        else:                   # arch down
            y_mid, y_edge = y1 - glyph_h / 2, y0 + glyph_h / 2
        # the rescue-fragment bbox is vertically inflated, so bbox-derived
        # anchors sit tens of px off the real glyph band (measured: mid
        # anchor 39px low, parabola depth 172px vs ~97px true on page 3).
        # Re-anchor on measured ink-run centers (text-colored rows) at the
        # chord center and +-0.65 flanks; two arc lines share the ribbon,
        # so each probe picks the run nearest its own prediction.
        if img is not None and block.style.text_rgb is not None:
            import numpy as np
            xc, half_w = (x0 + x1) / 2, (x1 - x0) / 2
            tc = np.asarray(block.style.text_rgb, dtype=int)

            def ink_center(t):
                y_pred = y_mid + (y_edge - y_mid) * t * t
                px = xc + t * half_w
                x_lo = max(0, int(px - glyph_h))
                x_hi = min(img.shape[1], int(px + glyph_h))
                y_lo = max(0, int(y_pred - 1.6 * glyph_h))
                y_hi = min(img.shape[0], int(y_pred + 1.6 * glyph_h))
                if x_hi - x_lo < 8 or y_hi - y_lo < 8:
                    return None
                win = img[y_lo:y_hi, x_lo:x_hi].astype(int)
                inkrow = (np.abs(win - tc).max(axis=2) < 40).sum(axis=1) >= 4
                runs, s, last = [], None, None
                for j, v in enumerate(inkrow):
                    if v:
                        if s is None:
                            if last is not None and runs and \
                                    j - last < 8:  # bridge AA breaks
                                s = runs.pop()[0]
                            else:
                                s = j
                        last = j
                    elif s is not None:
                        runs.append((s, last))
                        s = None
                if s is not None:
                    runs.append((s, last))
                runs = [r for r in runs if r[1] - r[0] >= 0.35 * glyph_h]
                if not runs:
                    return None
                centers = [y_lo + (a + b) / 2 for a, b in runs]
                return min(centers, key=lambda c: abs(c - y_pred))

            c0 = ink_center(0.0)
            if c0 is not None and abs(c0 - y_mid) < 1.2 * glyph_h:
                depths = [(yf - c0) / 0.4225
                          for yf in (ink_center(-0.65), ink_center(0.65))
                          if yf is not None]
                y_edge = c0 + (sum(depths) / len(depths) if depths
                               else y_edge - y_mid)
                y_mid = c0
        return ln, x0, y0, x1, y1, glyph_h, y_mid, y_edge

    @staticmethod
    def _ribbon_limits(img, sx0: float, sx1: float, yc: float, slope: float,
                       top: float, bot: float, fill_rgb,
                       glyph_h: float) -> tuple[float, float]:
        """Shrink a cover strip's vertical extent to the ribbon it sits on.

        Samples three columns across the strip and walks up/down from the
        local tangent line, following pixels near the fill color. Glyph
        strokes interrupt the ribbon only near the parabola center, so
        long gaps (up to glyph_h) are bridged inside a +-0.9*glyph_h zone;
        outside it only ~noise-sized gaps are allowed — otherwise the walk
        bridges across the ribbon's edge band and white gap onto the
        ribbon-colored dashes above and never clamps. Limits are taken in
        the tangent frame (min/max over the samples), so the rotated
        strip's corners stay inside too. Only ever shrinks — if the ribbon
        fills the whole window the limits come back unchanged."""
        import numpy as np

        h, w = img.shape[:2]
        fill = np.asarray(fill_rgb, dtype=int)
        zone = 0.9 * glyph_h
        gap_in, gap_out = int(glyph_h), 10
        cx_m = (sx0 + sx1) / 2
        r_top, r_bot = top - yc, bot - yc
        off_top, off_bot = [], []
        for fx in (0.15, 0.5, 0.85):
            cx = sx0 + (sx1 - sx0) * fx
            yci = yc + slope * (cx - cx_m)
            x_lo, x_hi = max(0, int(cx) - 2), min(w, int(cx) + 3)
            y_lo = max(0, int(yci + r_top))
            y_hi = min(h, int(yci + r_bot) + 2)
            if x_hi <= x_lo or y_hi - y_lo < 4:
                continue
            band = img[y_lo:y_hi, x_lo:x_hi].astype(int)
            col = (np.abs(band - fill).max(axis=2) < 32).mean(axis=1) >= 0.5

            def walk(idx_range):
                last, gap = None, 0
                for j in idx_range:
                    if col[j]:
                        last, gap = j, 0
                        continue
                    if last is None:
                        continue
                    gap += 1
                    allowed = gap_in if abs(y_lo + j - yci) <= zone \
                        else gap_out
                    if gap > allowed:
                        break
                return last

            start = min(max(int(round(yci)) - y_lo, 0), len(col) - 1)
            down = walk(range(start, len(col)))
            up = walk(range(start, -1, -1))
            if down is not None:
                off_bot.append(y_lo + down + 2 - yci)
            if up is not None:
                off_top.append(y_lo + up - 2 - yci)
        new_bot = min(bot, yc + min(off_bot)) if off_bot else bot
        new_top = max(top, yc + max(off_top)) if off_top else top
        if new_bot - new_top < 4:
            return top, bot
        return new_top, new_bot

    def _add_arc_cover(self, slide, block, i: int, ex, ey,
                       px_per_pt: float, img=None) -> None:
        """Cover strips tracing the arc band (drawn before ALL arc text)."""
        import numpy as np

        ln, x0, y0, x1, y1, glyph_h, y_mid, y_edge = \
            self._arc_geometry(block, px_per_pt, img)
        style = block.style

        if not (self.cover and style.bg_rgb is not None):
            return

        bg = np.asarray(style.bg_rgb, dtype=int)
        import math

        # ONE solid color for the whole arc (user-verified: a single
        # accurate mid-tone blends fine, while per-strip colors/gradients
        # introduce visible seams between strips): median of all
        # ribbon-family pixels under the arc band
        fill_rgb = style.bg_rgb
        if img is not None:
            region = img[max(0, int(y0)):max(1, int(y1)),
                         max(0, int(x0)):max(1, int(x1))]
            if region.size:
                px = region.reshape(-1, 3).astype(int)
                near = np.abs(px - bg).max(axis=1) < 70
                if near.sum() >= 100:
                    fill_rgb = tuple(int(v) for v in
                                     np.median(px[near], axis=0))

        # rotated strips at the local tangent: a horizontal rectangle on
        # the sloped flank pokes its corner past the ribbon edge no matter
        # how it is clamped; a tangent-aligned one hugs it
        n = 12
        xc, half_w = (x0 + x1) / 2, (x1 - x0) / 2
        # 1.5x the glyph height: enough slack (~0.25 glyph_h per side) for
        # parabola-fit residuals, small enough to leave the ribbon's soft
        # edge gradients untouched (2.2x flattened the top edge into a
        # hard line — user-reported)
        strip_h = glyph_h * 1.5
        # adjacent strips differ in rotation; rotating about their own
        # centers opens wedge gaps at the far corners (~strip_h/2 * dslope
        # per side) where the original raster peeks through — widen each
        # strip by the full wedge so neighbors overlap
        d_slope = abs(2 * (y_edge - y_mid) / half_w) * ((x1 - x0) / n) \
            / half_w
        extra_px = strip_h * d_slope
        for s in range(n):
            sx0 = x0 + (x1 - x0) * s / n
            sx1 = x0 + (x1 - x0) * (s + 1) / n
            cx_s = (sx0 + sx1) / 2
            t = (cx_s - xc) / half_w
            yc = y_mid + (y_edge - y_mid) * t * t
            slope = 2 * (y_edge - y_mid) * t / half_w
            ang = math.degrees(math.atan(slope))
            top, bot = yc - strip_h / 2, yc + strip_h / 2
            if img is not None:
                # the OCR bbox of an arc line is inflated by the rescue
                # fragments, which can push the strip bottom past the
                # ribbon's edge band — shrink (never grow) to the ribbon
                # extent actually present under the strip center
                top, bot = self._ribbon_limits(
                    img, sx0, sx1, yc, slope, top, bot, fill_rgb, glyph_h)
            width = Emu(round((ex(sx1) - ex(sx0))
                              / max(0.5, math.cos(math.radians(ang))) * 1.04
                              + ex(extra_px)))
            height = Emu(ey(bot) - ey(top))
            strip = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Emu(max(0, ex(cx_s) - width // 2)),
                Emu(max(0, ey((top + bot) / 2) - height // 2)),
                width, height,
            )
            strip.rotation = ang
            strip.name = f"Text {i} cover {s}"
            strip.shadow.inherit = False
            strip.line.fill.background()
            strip.fill.solid()
            strip.fill.fore_color.rgb = RGBColor(*fill_rgb)

    def _add_arc_segments(self, slide, block, i: int, ex, ey,
                          px_per_pt: float, img=None) -> None:
        """The arc text: chord segments along the parabola, each a small
        rotated shape at the local tangent. PowerPoint's prstTxWarp arch
        was tried first and abandoned — its drop/frame scaling is opaque
        and uncontrollable (see git history for the calibration attempts)."""
        import math

        ln, x0, y0, x1, y1, glyph_h, y_mid_t, y_edge_t = \
            self._arc_geometry(block, px_per_pt, img)
        style = block.style
        xc, half_w = (x0 + x1) / 2, (x1 - x0) / 2
        n_seg = max(3, min(6, round((x1 - x0) / (3.0 * max(glyph_h, 1)))))
        text = ln.text
        # split by accumulated advance width (mixed CJK/latin) and snap to
        # a nearby space so words like PDF are not cut in half
        adv = [1.0 if ord(c) >= 0x2E80 else (0.33 if c == " " else 0.52)
               for c in text]
        total = sum(adv) or 1.0
        bounds = [0]
        acc, target_idx = 0.0, 1
        for idx, a in enumerate(adv):
            acc += a
            while target_idx < n_seg and acc >= total * target_idx / n_seg:
                cut = idx + 1
                for off in (0, 1, -1, 2, -2):
                    j = cut + off
                    if 0 < j < len(text) and text[j - 1] == " ":
                        cut = j
                        break
                bounds.append(max(cut, bounds[-1]))
                target_idx += 1
        bounds.append(len(text))
        # pack segments by their actual text width, centered on the chord:
        # spreading them across the full chord turns the font-cap slack
        # into visible gaps between segments
        font_px = style.font_pt * px_per_pt
        em_at = [0.0]
        for s in range(n_seg):
            seg_em_s = sum(adv[bounds[s]:bounds[s + 1]])
            em_at.append(em_at[-1] + seg_em_s)
        em_total = em_at[-1] or 1.0
        span = min(x1 - x0, em_total * font_px * 1.0)
        sx_origin = xc - span / 2
        for s in range(n_seg):
            seg_text = text[bounds[s]:bounds[s + 1]].strip()
            if not seg_text:
                continue
            em_mid = (em_at[s] + em_at[s + 1]) / 2
            cx_seg = sx_origin + span * em_mid / em_total
            t = (cx_seg - xc) / half_w
            yc = y_mid_t + (y_edge_t - y_mid_t) * t * t
            slope = 2 * (y_edge_t - y_mid_t) * t / half_w
            ang = math.degrees(math.atan(slope))
            seg_em = sum(1.0 if ord(c) >= 0x2E80 else
                         (0.33 if c == " " else 0.52) for c in seg_text)
            width = Emu(round(seg_em * style.font_pt * EMU_PER_PT * 1.06))
            height = Emu(ey(yc + glyph_h / 2) - ey(yc - glyph_h / 2))
            left = Emu(ex(cx_seg) - width // 2)
            top = Emu(ey(yc) - height // 2)
            seg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top,
                                         width, height)
            seg.name = f"Text {i}.{s}"
            seg.rotation = ang
            seg.shadow.inherit = False
            seg.line.fill.background()
            seg.fill.background()
            tf = seg.text_frame
            tf.word_wrap = False
            tf.auto_size = None
            tf.margin_left = tf.margin_right = 0
            tf.margin_top = tf.margin_bottom = 0
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            run = para.add_run()
            run.text = seg_text
            font = run.font
            font.size = Pt(style.font_pt)
            font.bold = style.bold
            font.name = _latin_font(style.font_pt)
            font.color.rgb = RGBColor(*style.text_rgb)
            _set_east_asian_font(run, self.font_name)

    def save(self, path: str) -> None:
        self.prs.save(path)


def _mark_superscript(pieces, n_tail: int):
    """Tag the last n_tail characters of the pieces as superscript, splitting
    the piece they fall in. Returns [(text, rgb, is_super), ...]."""
    if n_tail <= 0:
        return [(t, c, False) for t, c in pieces]
    out, rem = [], n_tail
    for t, c in reversed(pieces):
        if rem <= 0:
            out.append((t, c, False))
        elif len(t) <= rem:
            out.append((t, c, True))
            rem -= len(t)
        else:
            out.append((t[len(t) - rem:], c, True))
            out.append((t[:len(t) - rem], c, False))
            rem = 0
    return [(t, c, s) for t, c, s in reversed(out) if t]


def _split_text_runs(text: str, runs) -> list[tuple[str, tuple]] | None:
    """Split the line text into (piece, rgb) runs. `runs` counts cover the
    space-stripped text; spaces and any trailing extras (recovered
    punctuation) attach to the current/last run."""
    total = sum(n for n, _ in runs)
    stripped_len = sum(1 for c in text if c != " ")
    if stripped_len < total:
        return None  # text was transformed past recognition; play safe

    bounds = []  # exclusive cumulative end index per run
    acc = 0
    for n, _ in runs:
        acc += n
        bounds.append(acc)

    pieces: list[tuple[str, tuple]] = []
    buf, seg, idx = [], 0, 0
    for ch in text:
        if ch != " ":
            while seg < len(runs) - 1 and idx >= bounds[seg]:
                pieces.append(("".join(buf), runs[seg][1]))
                buf = []
                seg += 1
            idx += 1
        buf.append(ch)
    if buf:
        pieces.append(("".join(buf), runs[seg][1]))
    return [(p, rgb) for p, rgb in pieces if p]


def _set_east_asian_font(run, typeface: str) -> None:
    """python-pptx has no API for <a:ea>; without it CJK text silently falls
    back to the theme font. Insert it right after <a:latin> (schema order)."""
    rPr = run._r.get_or_add_rPr()
    latin = rPr.find(qn("a:latin"))
    ea = rPr.makeelement(qn("a:ea"), {"typeface": typeface})
    if latin is not None:
        rPr.insert(list(rPr).index(latin) + 1, ea)
    else:
        rPr.append(ea)
