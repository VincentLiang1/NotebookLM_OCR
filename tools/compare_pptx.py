"""Compare a generated PPTX against a reference PPTX (e.g. the manually
corrected DeckEdit output): per-slide text recall via fuzzy matching.

Usage: python tools/compare_pptx.py generated.pptx reference.pptx
"""
from __future__ import annotations

import sys
import unicodedata
from difflib import SequenceMatcher

from pptx import Presentation


def slide_texts(prs) -> list[list[str]]:
    out = []
    for slide in prs.slides:
        texts = []
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            t = norm(sh.text_frame.text)
            if t:
                texts.append(t)
        out.append(texts)
    return out


def norm(s: str) -> str:
    s = unicodedata.normalize("NFKC", s)
    return "".join(s.split())  # drop all whitespace; CJK has none anyway


def best_match(needle: str, haystack: list[str]) -> tuple[float, str]:
    best, best_t = 0.0, ""
    for t in haystack:
        r = SequenceMatcher(None, needle, t).ratio()
        if r > best:
            best, best_t = r, t
    return best, best_t


def main() -> int:
    gen_path, ref_path = sys.argv[1], sys.argv[2]
    gen = slide_texts(Presentation(gen_path))
    ref = slide_texts(Presentation(ref_path))
    if len(gen) != len(ref):
        print(f"slide count differs: generated={len(gen)} reference={len(ref)}")

    total_ref = matched = 0
    for i, (g, r) in enumerate(zip(gen, ref), 1):
        hits = 0
        misses = []
        for t in r:
            score, _ = best_match(t, g)
            if score >= 0.85:
                hits += 1
            else:
                misses.append((score, t))
        total_ref += len(r)
        matched += hits
        line = f"slide {i:2d}: ref={len(r):2d} gen={len(g):2d} matched={hits:2d}"
        print(line)
        for score, t in misses:
            near = best_match(t, g)[1]
            print(f"    MISS ({score:.2f}) ref={t!r}")
            if near:
                print(f"               gen~{near!r}")
    print(f"\noverall recall: {matched}/{total_ref} "
          f"({100 * matched / max(1, total_ref):.1f}%)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
